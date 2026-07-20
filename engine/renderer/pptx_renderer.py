import json
import logging
from datetime import datetime, timezone
from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from engine.design.types import Theme

logger = logging.getLogger(__name__)

SLIDE_WIDTH  = Emu(9144000)
SLIDE_HEIGHT = Emu(5143500)


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, hex_color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _add_textbox(slide, left, top, width, height, text, font_name, font_size,
                  bold=False, italic=False, color="#000000",
                  align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txBox


class PPTXRenderer:
    def __init__(self, output_dir: Path, theme: Theme):
        self.output_dir = output_dir
        self.theme = theme

    def render(
        self,
        slides_json: dict,
        topic: str,
        provider_used: str,
        generation_time_ms: int = 0,
    ) -> tuple[str, Path]:
        job_id = str(uuid4())
        job_dir = self.output_dir / job_id
        job_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width  = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        blank_layout = prs.slide_layouts[6]

        handlers = {
            "title":      self._render_title,
            "bullets":    self._render_bullets,
            "two_column": self._render_two_column,
            "quote":      self._render_quote,
            "closing":    self._render_closing,
        }

        for slide_data in slides_json.get("slides", []):
            slide_type = slide_data.get("type", "bullets")
            handler = handlers.get(slide_type)
            if handler is None:
                logger.warning("미지원 타입 '%s' → bullets fallback", slide_type)
                handler = self._render_bullets
                slide_data = {
                    **slide_data, "type": "bullets",
                    "items": slide_data.get("items", [str(slide_data.get("quote", "내용"))]),
                }
            slide = prs.slides.add_slide(blank_layout)
            handler(slide, slide_data)
            notes = slide_data.get("speaker_notes", "")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        pptx_path = job_dir / "presentation.pptx"
        prs.save(pptx_path)

        metadata = {
            "job_id": job_id,
            "created_at": datetime.now(timezone.utc).isoformat(),
            "topic": topic,
            "provider_used": provider_used,
            "slide_count": len(slides_json.get("slides", [])),
            "generation_time_ms": generation_time_ms,
            "theme_name": self.theme.name,
        }
        (job_dir / "metadata.json").write_text(
            json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        logger.info("PPTX 생성 완료: %s (theme=%s)", pptx_path, self.theme.name)
        return job_id, pptx_path

    def _c(self, key: str) -> str:
        """컬러 단축 접근자"""
        return getattr(self.theme.colors, key)

    def _f(self, key: str):
        """폰트 크기 단축 접근자"""
        return getattr(self.theme.fonts.sizes, key)

    def _render_title(self, slide, data: dict):
        _set_bg(slide, self._c("primary"))
        w, h = SLIDE_WIDTH, SLIDE_HEIGHT
        _add_textbox(
            slide, Emu(int(w*0.05)), Emu(int(h*0.28)), Emu(int(w*0.90)), Emu(int(h*0.28)),
            data.get("title", ""), self.theme.fonts.title_font, self._f("cover_title"),
            bold=True, color=self._c("primary_text"), align=PP_ALIGN.CENTER,
        )
        subtitle = data.get("subtitle", "")
        if subtitle:
            _add_textbox(
                slide, Emu(int(w*0.05)), Emu(int(h*0.60)), Emu(int(w*0.90)), Emu(int(h*0.16)),
                subtitle, self.theme.fonts.body_font, self._f("cover_subtitle"),
                color=self._c("primary_text"), align=PP_ALIGN.CENTER,
            )

    def _render_bullets(self, slide, data: dict):
        _set_bg(slide, self._c("background"))
        w, h = SLIDE_WIDTH, SLIDE_HEIGHT
        bar = slide.shapes.add_shape(1, Emu(0), Emu(0), w, Emu(int(h*0.18)))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(self._c("primary")); bar.line.fill.background()
        _add_textbox(
            slide, Emu(int(w*0.04)), Emu(int(h*0.02)), Emu(int(w*0.92)), Emu(int(h*0.15)),
            data.get("title", ""), self.theme.fonts.title_font, self._f("slide_title"),
            bold=True, color=self._c("primary_text"),
        )
        items = data.get("items", [])
        txBox = slide.shapes.add_textbox(Emu(int(w*0.06)), Emu(int(h*0.22)), Emu(int(w*0.88)), Emu(int(h*0.70)))
        tf = txBox.text_frame; tf.word_wrap = True
        for i, item in enumerate(items[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run(); run.text = f"• {item}"
            run.font.name = self.theme.fonts.body_font
            run.font.size = Pt(self._f("bullet"))
            run.font.color.rgb = _rgb(self._c("text"))

    def _render_two_column(self, slide, data: dict):
        _set_bg(slide, self._c("background"))
        w, h = SLIDE_WIDTH, SLIDE_HEIGHT
        _add_textbox(
            slide, Emu(int(w*0.04)), Emu(int(h*0.04)), Emu(int(w*0.92)), Emu(int(h*0.14)),
            data.get("title", ""), self.theme.fonts.title_font, self._f("slide_title"),
            bold=True, color=self._c("text"),
        )
        col_top = Emu(int(h*0.20)); col_h = Emu(int(h*0.72)); col_w = Emu(int(w*0.44))
        bg_colors = [self._c("surface"), self._c("background")]
        for i, key in enumerate(["left", "right"]):
            col_left = Emu(int(w*(0.03 + i*0.50)))
            col_data = data.get(key, {})
            box = slide.shapes.add_shape(1, col_left, col_top, col_w, col_h)
            box.fill.solid(); box.fill.fore_color.rgb = _rgb(bg_colors[i]); box.line.color.rgb = _rgb("#E2E8F0")
            if heading := col_data.get("heading", ""):
                _add_textbox(
                    slide, col_left+Emu(100000), col_top+Emu(80000), col_w-Emu(200000), Emu(int(h*0.12)),
                    heading, self.theme.fonts.title_font, self._f("body"),
                    bold=True, color=self._c("primary"),
                )
            itb = slide.shapes.add_textbox(col_left+Emu(100000), col_top+Emu(int(h*0.14)), col_w-Emu(200000), col_h-Emu(int(h*0.16)))
            tf = itb.text_frame; tf.word_wrap = True
            for j, item in enumerate(col_data.get("items", [])):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                run = p.add_run(); run.text = f"• {item}"
                run.font.name = self.theme.fonts.body_font
                run.font.size = Pt(self._f("bullet"))
                run.font.color.rgb = _rgb(self._c("text"))

    def _render_quote(self, slide, data: dict):
        _set_bg(slide, self._c("primary"))
        w, h = SLIDE_WIDTH, SLIDE_HEIGHT
        _add_textbox(slide, Emu(int(w*0.05)), Emu(int(h*0.08)), Emu(int(w*0.15)), Emu(int(h*0.20)),
                     '"', self.theme.fonts.title_font, 80, color="#FFFFFF")
        _add_textbox(
            slide, Emu(int(w*0.10)), Emu(int(h*0.22)), Emu(int(w*0.80)), Emu(int(h*0.48)),
            data.get("quote", ""), self.theme.fonts.body_font, self._f("quote"),
            italic=True, color=self._c("primary_text"), align=PP_ALIGN.CENTER,
        )
        if source := data.get("source", ""):
            _add_textbox(
                slide, Emu(int(w*0.10)), Emu(int(h*0.74)), Emu(int(w*0.80)), Emu(int(h*0.10)),
                f"— {source}", self.theme.fonts.body_font, self._f("source"),
                color=self._c("primary_text"), align=PP_ALIGN.CENTER,
            )

    def _render_closing(self, slide, data: dict):
        _set_bg(slide, self._c("primary"))
        w, h = SLIDE_WIDTH, SLIDE_HEIGHT
        _add_textbox(
            slide, Emu(int(w*0.05)), Emu(int(h*0.30)), Emu(int(w*0.90)), Emu(int(h*0.25)),
            data.get("title", "감사합니다"), self.theme.fonts.title_font, self._f("cover_title"),
            bold=True, color=self._c("primary_text"), align=PP_ALIGN.CENTER,
        )
        if contact := data.get("contact", ""):
            _add_textbox(
                slide, Emu(int(w*0.05)), Emu(int(h*0.60)), Emu(int(w*0.90)), Emu(int(h*0.14)),
                contact, self.theme.fonts.body_font, self._f("cover_subtitle"),
                color=self._c("primary_text"), align=PP_ALIGN.CENTER,
            )
